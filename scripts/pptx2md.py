#!/usr/bin/env python3
"""Convert PowerPoint (.pptx) presentations to Markdown.

Features:
  - Slide titles -> Markdown headings (## per slide)
  - Tables -> HTML tables with rowspan/colspan for merged cells
  - Multiple text boxes sorted by visual position (top-to-bottom, left-to-right)
  - Inline formatting: bold, italic, underline, strikethrough, code, hyperlinks
  - Bullet lists detection
  - Grouped shapes (recursive traversal)
  - Speaker notes appended per slide
  - Heuristic title detection for PPTs without standard placeholders

Usage:
  python pptx2md.py input.pptx [-o output.md]

Can also be imported as a module:
  from pptx2md import convert
  convert("input.pptx", "output.md")

Requires: python-pptx (pip install python-pptx)
"""

import sys
import re
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Inline formatting helpers
# ---------------------------------------------------------------------------

MONOSPACE_FONTS = {
    "consolas", "courier", "courier new", "monospace",
    "source code pro", "fira code", "menlo", "cascadia code",
    "cascadia mono", "jetbrains mono", "lucida console",
}


def _escape_md(text: str) -> str:
    """Escape Markdown special characters."""
    if not text:
        return ""
    for ch in (r"\\", "`", "*", "_", "[", "]", "#", "~"):
        text = text.replace(ch, "\\" + ch)
    return text


def _run_to_inlines(r_el, is_html: bool = False) -> str:
    """Convert a single <a:r> run element to inline Markdown or HTML."""
    rPr = r_el.find(qn("a:rPr"))
    is_bold = is_italic = is_underline = is_strike = is_code = False
    if rPr is not None:
        # In PPTX XML, formatting flags are ATTRIBUTES on <a:rPr>, not child elements
        if rPr.get("b") == "1":
            is_bold = True
        if rPr.get("i") == "1":
            is_italic = True
        u_val = rPr.get("u", "")
        if u_val and u_val != "none":
            is_underline = True
        s_val = rPr.get("strike", "")
        if s_val and s_val != "noStrike":
            is_strike = True
        # Monospace font detection (these ARE child elements)
        latin = rPr.find(qn("a:latin"))
        ea = rPr.find(qn("a:ea"))
        for font_el in (latin, ea):
            if font_el is not None:
                tf = (font_el.get("typeface") or "").lower()
                if tf in MONOSPACE_FONTS:
                    is_code = True
                    break

    t_el = r_el.find(qn("a:t"))
    text = t_el.text if t_el is not None and t_el.text else ""
    if not text:
        return ""

    if is_html:
        # HTML mode for table cells — wrap with HTML tags
        result = text
        if is_code:
            result = f"<code>{result}</code>"
        if is_strike:
            result = f"<s>{result}</s>"
        if is_underline:
            result = f"<u>{result}</u>"
        if is_italic:
            result = f"<em>{result}</em>"
        if is_bold:
            result = f"<strong>{result}</strong>"
        return result
    else:
        # Markdown mode for text boxes
        result = _escape_md(text)
        if not result.strip():
            return result  # whitespace-only, return as-is
        if is_code:
            return f"`{text}`"  # code blocks don't need inner escaping
        if is_strike:
            result = f"~~{result}~~"
        if is_underline:
            result = f"<u>{result}</u>"
        if is_italic:
            result = f"*{result}*"
        if is_bold:
            result = f"**{result}**"
        return result


def _para_to_inlines(p_el, is_html: bool = False) -> str:
    """Convert an <a:p> element to inline text, handling runs, line breaks, and fields."""
    parts = []
    for child in p_el:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "r":
            parts.append(_run_to_inlines(child, is_html))

        elif tag == "br":
            parts.append("<br>" if is_html else "  \n")

        elif tag == "fld":
            # Field (date, time, slide number, etc.)
            t_el = child.find(qn("a:t"))
            if t_el is not None and t_el.text:
                parts.append(t_el.text)

    return "".join(parts)


# ---------------------------------------------------------------------------
# Table cell text (with inline formatting in HTML mode)
# ---------------------------------------------------------------------------

def _cell_text_html(tc) -> str:
    """Get text content of a table cell as HTML (for use in HTML tables)."""
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return ""
    lines = []
    for p in txBody.findall(qn("a:p")):
        text = _para_to_inlines(p, is_html=True).strip()
        if text:
            lines.append(text)
    return "<br>".join(lines)


# ---------------------------------------------------------------------------
# Table -> HTML with merged-cell support
# ---------------------------------------------------------------------------

def _table_to_html(table) -> str:
    """Convert a PPTX table to an HTML table with rowspan/colspan.

    Reads raw XML <a:tr> / <a:tc> elements to correctly handle merged cells.
    In PPTX XML, merged cells use gridSpan/rowSpan attributes on the leading
    <a:tc>, and hMerge/vMerge continuation cells fill the remaining grid
    positions. The total <a:tc> count per <a:tr> always equals the grid width.
    """
    tbl = table._tbl  # raw <a:tbl> element
    trs = tbl.findall(qn("a:tr"))
    if not trs:
        return ""

    n_rows = len(trs)

    # Parse all rows: each row is a list of <a:tc> elements
    rows_tcs = []
    for tr in trs:
        tcs = tr.findall(qn("a:tc"))
        rows_tcs.append(tcs)

    # Determine grid width from all rows (take maximum)
    # Some rows may have fewer <a:tc> due to merges
    n_cols = 0
    for tcs in rows_tcs:
        row_width = 0
        for tc in tcs:
            gs = int(tc.get("gridSpan", "1"))
            row_width += gs
        n_cols = max(n_cols, row_width)
    if n_cols == 0:
        n_cols = max(len(r) for r in rows_tcs) if rows_tcs else 0
    if n_cols == 0:
        return ""

    # Track which (row, col) grid positions are already covered by a span
    covered = set()

    html = ["<table>"]

    for r, tcs in enumerate(rows_tcs):
        html.append("  <tr>")
        col = 0  # grid column tracker

        for tc in tcs:
            gs = int(tc.get("gridSpan", "1"))

            # Check for hMerge / vMerge continuation cells
            h_merge = tc.get("hMerge")
            v_merge = tc.get("vMerge")

            if h_merge == "1" or v_merge == "1":
                # Continuation of a merged cell — the leading cell's
                # gridSpan/rowSpan already covers this position.
                # Do NOT advance col; just skip this <a:tc>.
                continue

            # If this grid position is already covered by a previous span
            # (e.g. rowspan from a prior row), skip it
            while col < n_cols and (r, col) in covered:
                col += 1
            if col >= n_cols:
                break

            # Read rowspan directly from rowSpan attribute
            rs = int(tc.get("rowSpan", "1"))

            # Mark covered positions
            for dr in range(rs):
                for dc in range(gs):
                    if dr == 0 and dc == 0:
                        continue
                    covered.add((r + dr, col + dc))

            # Build cell HTML
            text = _cell_text_html(tc)
            attrs = ""
            if gs > 1:
                attrs += f' colspan="{gs}"'
            if rs > 1:
                attrs += f' rowspan="{rs}"'

            tag = "th" if r == 0 else "td"
            html.append(f"    <{tag}{attrs}>{text}</{tag}>")
            col += gs

        html.append("  </tr>")

    html.append("</table>")
    return "\n".join(html)


# ---------------------------------------------------------------------------
# Shape processing (text boxes, groups, tables)
# ---------------------------------------------------------------------------

def _is_title_placeholder(shape) -> bool:
    """Check if a shape is a title/subtitle placeholder."""
    if not shape.is_placeholder:
        return False
    ph_idx = shape.placeholder_format.idx
    # Title=0 or 1, Center Title=2, Subtitle=3
    return ph_idx in (0, 1, 2, 3)


def _title_priority(shape) -> int:
    """Return title priority (lower = higher priority). Used to pick the main title."""
    if not shape.is_placeholder:
        return 99
    ph_idx = shape.placeholder_format.idx
    # Priority: Title(0) > Center Title(2) > Subtitle(1,3)
    return {0: 0, 2: 1, 1: 2, 3: 3}.get(ph_idx, 50)


_NUMERIC_RE = re.compile(
    r'^[\s]*[-+]?[\s]*[\d,]+\.?\d*[\s]*'
    r'(万元|元|亿|%|人天|个|张|条|家|项|月|年|日)?[\s]*$'
)


def _find_heuristic_title(slide, slide_height: int) -> tuple:
    """Find a title-like shape using heuristics when no placeholder title exists.

    Strategy: among shapes in the top half with short, non-numeric text,
    pick the one with the largest font size (checking both rPr and defRPr).
    Ties broken by highest position (smallest top).
    Fallback: if nothing found in top half, search entire slide.

    Returns (title_text, shape) or ("", None).
    """
    MAX_TEXT_LEN = 60
    TOP_THRESHOLD = slide_height * 0.50  # top half of slide

    def _get_font_size(shape) -> int:
        """Get font size in EMU from first paragraph's first run or defRPr."""
        paras = shape.text_frame.paragraphs
        if not paras:
            return 0
        p_el = paras[0]._p
        runs = p_el.findall(qn("a:r"))
        sz = None
        if runs:
            rPr = runs[0].find(qn("a:rPr"))
            if rPr is not None:
                sz = rPr.get("sz")
        if sz is None:
            pPr = p_el.find(qn("a:pPr"))
            if pPr is not None:
                defRPr = pPr.find(qn("a:defRPr"))
                if defRPr is not None:
                    sz = defRPr.get("sz")
        if sz is None:
            return 0
        return int(sz) * 12700  # hundredths of a point -> EMU

    def _is_numeric(text: str) -> bool:
        """Check if text is a data point / statistic rather than a real title.

        Catches pure numbers ("123", "1,000") as well as short text with a
        high proportion of Arabic digits ("≈2000+条", "8月-9月", "47+").
        """
        # Pure numeric pattern (original check)
        if _NUMERIC_RE.match(text):
            return True
        # Digit-ratio heuristic: short text where ≥30 % of chars are digits
        if len(text) <= 20:
            digit_count = sum(1 for c in text if c.isdigit())
            if digit_count > 0 and digit_count / len(text) >= 0.3:
                return True
        return False

    def _collect_candidates(top_limit: float) -> list:
        candidates = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            if not shape.has_text_frame:
                continue
            if shape.top is None:
                continue

            text = shape.text_frame.text.strip()
            if not text or len(text) >= MAX_TEXT_LEN:
                continue

            if shape.top > top_limit:
                continue

            if _is_numeric(text):
                continue

            font_size = _get_font_size(shape)
            if font_size <= 0:
                continue

            candidates.append((shape, text, font_size, shape.top))
        return candidates

    # First try top half
    candidates = _collect_candidates(TOP_THRESHOLD)

    # Fallback: entire slide (for cover pages with title at bottom)
    if not candidates:
        candidates = _collect_candidates(slide_height)

    if not candidates:
        return ("", None)

    # Pick best: largest font, then highest position.
    # However, when multiple candidates share the max font size they are
    # likely infographic data-points (e.g. "7个", "≈2000+条") rather than
    # the real title.  In that case drop to the next font-size tier.
    max_sz = max(c[2] for c in candidates)
    max_count = sum(1 for c in candidates if c[2] == max_sz)
    if max_count > 1:
        # Keep only candidates with a strictly smaller font size
        smaller = [c for c in candidates if c[2] < max_sz]
        if smaller:
            candidates = smaller
        # else: all candidates share the same size — keep them all

    candidates.sort(key=lambda c: (-c[2], c[3]))
    best = candidates[0]
    return (best[1], best[0])


def _get_bullet_prefix(p_el, level: int) -> str:
    """Determine bullet marker for a paragraph. Returns '' if not a bullet."""
    pPr = p_el.find(qn("a:pPr"))
    if pPr is not None:
        # Explicit buNone -> not a bullet
        if pPr.find(qn("a:buNone")) is not None:
            return ""
        # Has explicit bullet markers
        if (pPr.find(qn("a:buChar")) is not None or
                pPr.find(qn("a:buAutoNum")) is not None):
            return "- "
    # Level > 0 usually implies bullet in PPTX
    if level > 0:
        return "- "
    return ""


def _get_auto_num(p_el, level: int) -> str | None:
    """If paragraph has auto-numbering, return the marker format."""
    pPr = p_el.find(qn("a:pPr"))
    if pPr is None:
        return None
    buAutoNum = pPr.find(qn("a:buAutoNum"))
    if buAutoNum is None:
        return None
    num_type = buAutoNum.get("type", "")
    # Types: arabicPeriod (1. 2. 3.), arabicParenBoth ((1) (2) (3)), etc.
    return num_type


def _tf_to_blocks(text_frame) -> list[str]:
    """Convert a text frame to a list of Markdown text blocks.

    Each block is a paragraph or list item. Blocks are separated by blank lines.
    """
    blocks = []
    auto_counter = 0

    for para in text_frame.paragraphs:
        p_el = para._p
        text = _para_to_inlines(p_el, is_html=False)
        level = para.level

        # Check for auto-numbering
        auto_type = _get_auto_num(p_el, level)
        if auto_type:
            auto_counter += 1
            if "Paren" in auto_type:
                marker = f"({auto_counter})"
            else:
                marker = f"{auto_counter}."
            indent = "  " * level
            blocks.append(f"{indent}{marker} {text}")
            continue
        else:
            # Reset counter when hitting a non-numbered paragraph
            if not text.strip():
                continue
            auto_counter = 0

        # Check for bullet
        bullet = _get_bullet_prefix(p_el, level)
        if bullet:
            indent = "  " * level
            blocks.append(f"{indent}{bullet}{text}")
        else:
            if text.strip():
                blocks.append(text)

    return blocks


def _process_shape(shape, content_parts: list[str], tables_seen: set):
    """Process a single shape and append content blocks."""
    # Group shape: recurse into children
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            _process_shape(s, content_parts, tables_seen)
        return

    # Table
    if shape.has_table:
        tables_seen.add(id(shape))
        content_parts.append("")
        content_parts.append(_table_to_html(shape.table))
        content_parts.append("")
        return

    # Text content
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if not text:
            return
        blocks = _tf_to_blocks(shape.text_frame)
        if blocks:
            content_parts.append("")
            content_parts.extend(blocks)
            content_parts.append("")


def _sorted_content_shapes(slide):
    """Return non-title shapes sorted by visual position (top-to-bottom, left-to-right)."""
    shapes = []
    for shape in slide.shapes:
        if _is_title_placeholder(shape):
            continue
        # Skip shapes with no visible content
        if not (shape.has_text_frame or shape.has_table or
                shape.shape_type == MSO_SHAPE_TYPE.GROUP):
            continue
        shapes.append(shape)

    # Sort by top position, then left position
    def sort_key(s):
        t = s.top if s.top is not None else 0
        l = s.left if s.left is not None else 0
        return (t, l)

    shapes.sort(key=sort_key)
    return shapes


# ---------------------------------------------------------------------------
# Slide processing
# ---------------------------------------------------------------------------

def _process_slide(slide, slide_num: int, total: int, slide_height: int = 6858000) -> str:
    """Convert a single slide to Markdown."""
    parts = []

    # --- Title: pick the highest-priority title placeholder ---
    title_text = ""
    title_shape = None          # strong ref prevents id() collision via GC
    best_priority = 99
    subtitle_parts = []

    for shape in slide.shapes:
        if _is_title_placeholder(shape) and shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if not t:
                continue
            prio = _title_priority(shape)
            if prio < best_priority:
                # Demote previous title to subtitle
                if title_text:
                    subtitle_parts.append(title_text)
                title_text = t
                title_shape = shape
                best_priority = prio
            else:
                # This is a lower-priority title (e.g. subtitle)
                subtitle_parts.append(t)

    # Fallback: heuristic title detection for PPTs without placeholders
    if not title_text:
        title_text, title_shape = _find_heuristic_title(slide, slide_height)

    # --- Content shapes (sorted by position) ---
    content_parts = []
    tables_seen = set()

    for shape in _sorted_content_shapes(slide):
        if shape is title_shape:
            continue
        # Skip other title/subtitle placeholders (already handled)
        if _is_title_placeholder(shape):
            continue
        _process_shape(shape, content_parts, tables_seen)

    # --- Speaker notes ---
    notes_text = ""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text.strip()

    # --- Assemble ---
    heading = f"## {title_text}" if title_text else f"## Slide {slide_num}"
    parts.append(heading)

    # Add subtitle text right after heading
    if subtitle_parts:
        for sp in subtitle_parts:
            parts.append(sp)

    content = "\n".join(content_parts).strip()
    if content:
        parts.append(content)

    if notes_text:
        parts.append("")
        parts.append(f"> **Notes:** {notes_text}")

    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Main conversion
# ---------------------------------------------------------------------------

def convert(pptx_path: str, output_path: str | None = None) -> str:
    """Convert a .pptx file to Markdown and return the Markdown string."""
    prs = Presentation(pptx_path)
    slides_md = []
    slide_height = prs.slide_height or 6858000  # default 10 inches in EMU

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        md = _process_slide(slide, i, total, slide_height)
        slides_md.append(md)

    result = "\n\n---\n\n".join(slides_md)

    # Clean up excessive blank lines
    while "\n\n\n\n" in result:
        result = result.replace("\n\n\n\n", "\n\n\n")
    while "\n\n\n" in result:
        result = result.replace("\n\n\n", "\n\n")
    result = result.strip() + "\n"

    if output_path:
        Path(output_path).write_text(result, encoding="utf-8")
        print(f"OK: {pptx_path} -> {output_path}  ({total} slides)", file=sys.stderr)

    return result


def main():
    ap = argparse.ArgumentParser(
        description="Convert PowerPoint (.pptx) to Markdown",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    ap.add_argument("input", help="Input .pptx file path")
    ap.add_argument("-o", "--output", help="Output .md file (default: <input>.md)")
    args = ap.parse_args()

    inp = Path(args.input)
    if not inp.exists():
        print(f"Error: file not found: {inp}", file=sys.stderr)
        sys.exit(1)

    out = args.output or str(inp.with_suffix(".md"))
    convert(str(inp), out)


if __name__ == "__main__":
    main()
